import os
from datetime import datetime,timezone
from typing import Optional, List, Dict, Any
from fastapi import FastAPI, HTTPException,Request,Query
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field, BeforeValidator
from pymongo import MongoClient
from bson import ObjectId
from dotenv import load_dotenv
from typing_extensions import Annotated
import uuid
import base64
import json
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
import matplotlib.pyplot as plt
import networkx as nx
import time
from openpyxl import Workbook
from docx import Document
from pptx import Presentation
from fastapi.background import BackgroundTasks
from fastapi.responses import FileResponse, JSONResponse
from googletrans import Translator

# -------------------- ENV SETUP --------------------
load_dotenv()

MONGODB_URI = os.getenv("MONGODB_URI")
DB_NAME = os.getenv("DB_NAME", "margdarshak")

if not MONGODB_URI:
    raise RuntimeError("MONGODB_URI not set in environment")

# -------------------- DB CONNECTION --------------------
client = MongoClient(MONGODB_URI)
db = client[DB_NAME]

# -------------------- DB COLLECTION CREATION --------------------
organization_profiles_collection = db.organization_profiles
state_context_rules_collection = db.state_context_rules
ai_context_analysis_collection = db.ai_context_analysis
problem_statements_collection = db.problem_statements
problem_refinements_collection = db.problem_refinements
problem_tree_collection = db["problem_trees"]
ecosystem_patterns_collection = db["ecosystem_problem_patterns"]
district_challenges_collection = db["district_challenges"]
student_outcomes_collection = db["student_outcomes"]
competency_frameworks_collection = db["competency_frameworks"]
policy_references_collection = db["policy_references"]
methodology_library_collection = db["methodology_library"]
selected_methodologies_collection = db["selected_methodologies"]
theory_of_change_collection = db["theory_of_change"]
toc_pattern_library_collection = db["toc_pattern_library"]
stakeholder_master_collection = db["stakeholder_master"]
organization_stakeholders_collection = db["organization_stakeholders"]
practice_master_collection = db["practice_master"]
practice_change_collection = db["practice_changes"]
indicator_master_collection = db["indicator_master"]
generated_indicators_collection = db["generated_indicators"]
indicator_targets_collection = db["indicator_targets"]
lfa_completeness_collection = db["lfa_completeness_scores"]
design_quality_feedback_collection = db["design_quality_feedback"]
lfa_templates_collection = db["lfa_templates"]
organization_templates_collection = db["organization_templates"]
template_ratings_collection = db["template_ratings"]
export_jobs_collection = db["export_jobs"]

# -------------------- FASTAPI APP --------------------
app = FastAPI(
    title="MargDarshak Program Design Platform",
    version="1.0.0",
    description="AI-powered program design backend for education NGOs"
)

# -------------------- CORS --------------------
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# -------------------- UTILITIES --------------------
PyObjectId = Annotated[
    str,
    BeforeValidator(
        lambda v: str(v) if ObjectId.is_valid(v) else ValueError("Invalid ObjectId")
    )
]

def utc_now():
    return datetime.now(timezone.utc)

def serialize_mongo(doc):
    doc["_id"] = str(doc["_id"])
    return doc

# -------------------- ALL FUNCTIONALITIES --------------------

# -------------------- Organization Profile Builder --------------------
class Geography(BaseModel):
    state: str
    district: Optional[str] = None
    block: Optional[str] = None

class ReachMetrics(BaseModel):
    schools: Optional[int] = Field(None, ge=0)
    students: Optional[int] = Field(None, ge=0)
    teachers: Optional[int] = Field(None, ge=0)

class TeamExpertise(BaseModel):
    role: str
    count: int = Field(..., ge=0)

class OrganizationProfileCreate(BaseModel):
    organization_name: str
    geography: Geography
    thematic_focus: List[str]
    maturity_level: str
    reach_metrics: Optional[ReachMetrics]
    team_size: int = Field(..., ge=1)
    team_expertise: Optional[List[TeamExpertise]]

class OrganizationProfileDB(BaseModel):
    id: PyObjectId = Field(default_factory=PyObjectId, alias="_id")
    organization_name: str
    geography: Geography
    thematic_focus: List[str]
    maturity_level: str
    reach_metrics: Optional[ReachMetrics]
    team_size: int
    team_expertise: Optional[List[TeamExpertise]]
    created_at: datetime
    updated_at: datetime

    class Config:
        allow_population_by_field_name = True
        json_encoders = {ObjectId: str}

# Organization Profile Creation Endpoint
@app.post("/organization/profile", response_model=OrganizationProfileDB)
def create_organization_profile(payload: OrganizationProfileCreate):

    existing = organization_profiles_collection.find_one({
        "organization_name": payload.organization_name
    })

    if existing:
        raise HTTPException(
            status_code=409,
            detail="Organization profile already exists"
        )

    doc = payload.dict()
    doc["created_at"] = utc_now()
    doc["updated_at"] = utc_now()

    result = organization_profiles_collection.insert_one(doc)
    doc["_id"] = result.inserted_id

    return serialize_mongo(doc)

# Organization Profile Retrieval Endpoint
@app.get("/organization/profile/{org_id}", response_model=OrganizationProfileDB)
def get_organization_profile(org_id: str):

    if not ObjectId.is_valid(org_id):
        raise HTTPException(status_code=400, detail="Invalid organization ID")

    org = organization_profiles_collection.find_one(
        {"_id": ObjectId(org_id)}
    )

    if not org:
        raise HTTPException(status_code=404, detail="Organization not found")

    return serialize_mongo(org)

# -------------------- AI-Powered Context Analysis --------------------
# AI CONTEXT MODELS
class LFARecommendation(BaseModel):
    template_key: str
    rationale: str

class ProgramPatternSuggestion(BaseModel):
    pattern_name: str
    relevance_reason: str

class ContextChallenge(BaseModel):
    challenge: str
    reason: str

class AIContextAnalysisResponse(BaseModel):
    organization_id: str
    lfa_recommendation: LFARecommendation
    similar_program_patterns: List[ProgramPatternSuggestion]
    potential_challenges: List[ContextChallenge]
    generated_at: datetime

# AI CONTEXT ENGINE
def analyze_context(profile: Dict[str, Any]) -> Dict[str, Any]:

    geography = profile["geography"]["state"].lower()
    themes = [t.lower() for t in profile["thematic_focus"]]
    maturity = profile["maturity_level"].lower()

    # ---- LFA TEMPLATE RECOMMENDATION ----
    if "fln" in themes:
        lfa_template = {
            "template_key": "FLN_System_Strengthening",
            "rationale": "FLN requires system-wide literacy improvement across grades and teachers"
        }
    elif "career readiness" in themes:
        lfa_template = {
            "template_key": "Youth_Career_Pathways",
            "rationale": "Career readiness needs multi-actor coordination and long-term outcomes"
        }
    else:
        lfa_template = {
            "template_key": "Education_Improvement_Generic",
            "rationale": "Suitable for multi-theme education interventions"
        }

    # ---- PROGRAM PATTERNS ----
    patterns = []

    if "fln" in themes:
        patterns.append({
            "pattern_name": "Teacher Coaching + Classroom Observation",
            "relevance_reason": "Improves instructional quality and student literacy outcomes"
        })

    if maturity == "startup":
        patterns.append({
            "pattern_name": "Pilot → Iterate → Scale",
            "relevance_reason": "Minimizes risk and improves learning before expansion"
        })

    # ---- STATE-SPECIFIC CHALLENGES (DB-DRIVEN) ----
    challenges = []

    state_rules = state_context_rules_collection.find_one({"state": geography})

    if state_rules:
        for item in state_rules.get("education_challenges", []):
            challenges.append({
                "challenge": item["challenge"],
                "reason": item["reason"]
            })

    return {
        "lfa_recommendation": lfa_template,
        "similar_program_patterns": patterns,
        "potential_challenges": challenges
    }

# AI CONTEXT ANALYSIS ENDPOINT
@app.post("/organization/{org_id}/ai-context", response_model=AIContextAnalysisResponse)
def generate_ai_context(org_id: str):

    if not ObjectId.is_valid(org_id):
        raise HTTPException(status_code=400, detail="Invalid organization ID")

    profile = organization_profiles_collection.find_one({"_id": ObjectId(org_id)})

    if not profile:
        raise HTTPException(status_code=404, detail="Organization not found")

    analysis = analyze_context(profile)

    response_doc = {
        "organization_id": str(profile["_id"]),
        "lfa_recommendation": analysis["lfa_recommendation"],
        "similar_program_patterns": analysis["similar_program_patterns"],
        "potential_challenges": analysis["potential_challenges"],
        "generated_at": utc_now()
    }

    ai_context_analysis_collection.insert_one(response_doc)

    return response_doc

# -------------------- PROBLEM STATEMENT Creation --------------------
class ProblemEvidence(BaseModel):
    source: Optional[str] = Field(None, example="ASER 2023")
    description: str = Field(..., example="Only 42% of Grade 3 students can read a Grade 2 text")
    year: Optional[int] = Field(None, ge=2000)

class ProblemStatementCreate(BaseModel):
    organization_id: str
    core_problem: str = Field(..., min_length=20)
    affected_stakeholders: List[str] = Field(
        ..., example=["Students", "Teachers"]
    )
    evidence: Optional[List[ProblemEvidence]]

class ProblemStatementDB(BaseModel):
    id: PyObjectId = Field(default_factory=PyObjectId, alias="_id")
    organization_id: str
    core_problem: str
    affected_stakeholders: List[str]
    evidence: Optional[List[ProblemEvidence]]
    created_at: datetime
    updated_at: datetime

    class Config:
        allow_population_by_field_name = True
        json_encoders = {ObjectId: str}

# Problem Statement Creation Endpoint
@app.post("/problem-statement", response_model=ProblemStatementDB)
def create_problem_statement(payload: ProblemStatementCreate):

    # Validate organization exists
    if not ObjectId.is_valid(payload.organization_id):
        raise HTTPException(status_code=400, detail="Invalid organization ID")

    org = organization_profiles_collection.find_one(
        {"_id": ObjectId(payload.organization_id)}
    )

    if not org:
        raise HTTPException(status_code=404, detail="Organization not found")

    doc = payload.dict()
    doc["created_at"] = utc_now()
    doc["updated_at"] = utc_now()

    result = problem_statements_collection.insert_one(doc)
    doc["_id"] = result.inserted_id

    return doc

# Problem Statement Retrieval Endpoint
@app.get("/organization/{org_id}/problem-statements", response_model=List[ProblemStatementDB])
def get_problem_statements(org_id: str):

    if not ObjectId.is_valid(org_id):
        raise HTTPException(status_code=400, detail="Invalid organization ID")

    statements = list(
        problem_statements_collection.find(
            {"organization_id": org_id}
        )
    )

    return statements

# --------------------AI Problem Refinement Assistant --------------------
class ProblemIssue(BaseModel):
    issue_type: str  # vague, solution-biased, missing-actor, etc.
    description: str

class RootCauseSuggestion(BaseModel):
    cause: str
    rationale: str

class ProblemRefinementResponse(BaseModel):
    problem_statement_id: str
    clarity_score: int = Field(..., ge=1, le=5)
    identified_issues: List[ProblemIssue]
    refined_problem_statement: str
    suggested_root_causes: List[RootCauseSuggestion]
    generated_at: datetime

def refine_problem_statement(problem_text: str) -> Dict[str, Any]:

    text = problem_text.lower()

    issues = []
    root_causes = []
    refined_text = problem_text
    clarity_score = 5

    # 1. CLARITY & LANGUAGE ANALYSIS
    vague_terms = {
        "lack of": "Specify what is lacking and to what extent",
        "poor": "Quantify or describe the quality gap",
        "low": "Clarify baseline or benchmark",
        "inadequate": "State what standards are not being met",
        "insufficient": "Mention compared to what requirement"
    }

    solution_terms = {
        "train": "Avoid prescribing solutions in the problem statement",
        "provide": "Focus on the problem, not the intervention",
        "implement": "Problem statements should be solution-neutral",
        "introduce": "Describe the gap, not the action",
        "conduct": "Actions belong in intervention design"
    }

    outcome_terms = ["learning outcomes", "achievement", "literacy", "numeracy"]
    system_terms = ["system", "policy", "governance", "administration"]

    # Detect vague language
    for term, explanation in vague_terms.items():
        if term in text:
            issues.append({
                "issue_type": "vague_language",
                "description": f"Uses vague term '{term}'. {explanation}."
            })
            clarity_score -= 1

    # Detect solution bias
    for term, explanation in solution_terms.items():
        if term in text:
            issues.append({
                "issue_type": "solution_bias",
                "description": f"Contains solution-oriented term '{term}'. {explanation}."
            })
            clarity_score -= 1

    # Detect missing outcome focus
    if not any(term in text for term in outcome_terms):
        issues.append({
            "issue_type": "missing_outcome_focus",
            "description": "Problem does not clearly describe the student-level outcome being affected"
        })
        clarity_score -= 1

    # Detect unclear level (system vs classroom)
    if not any(term in text for term in system_terms) and "students" not in text:
        issues.append({
            "issue_type": "unclear_system_level",
            "description": "Problem does not specify whether the issue is at classroom, school, or system level"
        })
        clarity_score -= 1

    clarity_score = max(1, clarity_score)

    # 2. REFINED PROBLEM STATEMENT (NEUTRAL & PRECISE)
    refinement_map = {
        "lack of": "limited availability of",
        "poor": "suboptimal",
        "low": "below expected levels of",
        "inadequate": "not aligned with required standards",
        "insufficient": "not meeting the required threshold"
    }

    for term, replacement in refinement_map.items():
        refined_text = refined_text.replace(term, replacement)

    # 3. ROOT CAUSE INFERENCE (MULTI-LAYERED)
    # Student-related inference
    if "students" in text or "children" in text:
        root_causes.append({
            "cause": "Foundational skill gaps from earlier grades",
            "rationale": "Learning deficits often accumulate due to weak early-grade instruction"
        })

        root_causes.append({
            "cause": "Limited opportunities for practice and reinforcement",
            "rationale": "Students require repeated exposure and feedback to master skills"
        })

    # Teacher-related inference
    if "teachers" in text or "teaching" in text:
        root_causes.append({
            "cause": "Inconsistent instructional practices",
            "rationale": "Variation in pedagogy leads to uneven learning outcomes"
        })

        root_causes.append({
            "cause": "Limited ongoing academic support for teachers",
            "rationale": "Teachers often lack continuous coaching and feedback mechanisms"
        })

    # School / leadership inference
    if "school" in text or "headmaster" in text or "leadership" in text:
        root_causes.append({
            "cause": "Weak instructional leadership at school level",
            "rationale": "School leaders play a critical role in setting academic priorities"
        })

    # System-level inference
    if any(term in text for term in system_terms):
        root_causes.append({
            "cause": "Fragmented implementation across system layers",
            "rationale": "Misalignment between policy, administration, and classrooms reduces effectiveness"
        })

        root_causes.append({
            "cause": "Monitoring focused on compliance rather than learning",
            "rationale": "Systems often track inputs instead of learning quality and outcomes"
        })

    # Fallback if nothing triggered
    if not root_causes:
        root_causes.append({
            "cause": "Multi-level coordination gaps",
            "rationale": "Education challenges often arise from weak alignment across stakeholders"
        })

    # 4. RETURN STRUCTURED OUTPUT
    return {
        "clarity_score": clarity_score,
        "identified_issues": issues,
        "refined_problem_statement": refined_text,
        "suggested_root_causes": root_causes
    }

# AI PROBLEM REFINEMENT API
@app.post(
    "/problem-statement/{problem_id}/ai-refine",
    response_model=ProblemRefinementResponse
)
def ai_refine_problem_statement(problem_id: str):

    if not ObjectId.is_valid(problem_id):
        raise HTTPException(status_code=400, detail="Invalid problem statement ID")

    problem = problem_statements_collection.find_one(
        {"_id": ObjectId(problem_id)}
    )

    if not problem:
        raise HTTPException(status_code=404, detail="Problem statement not found")

    analysis = refine_problem_statement(problem["core_problem"])

    response_doc = {
        "problem_statement_id": str(problem["_id"]),
        "clarity_score": analysis["clarity_score"],
        "identified_issues": analysis["identified_issues"],
        "refined_problem_statement": analysis["refined_problem_statement"],
        "suggested_root_causes": analysis["suggested_root_causes"],
        "generated_at": utc_now()
    }

    problem_refinements_collection.insert_one(response_doc)

    return response_doc

# -------------------- Problem Tree Builder --------------------
class ProblemTreeRequest(BaseModel):
    organization_id: str
    state: str
    district: str
    theme: str
    refined_problem_statement: str
    suggested_root_causes: List[Dict[str, str]]

class ProblemTreeResponse(BaseModel):
    problem_tree: Dict[str, Any]
    mermaid_diagram: str
    mermaid_preview_url: str
    mermaid_png_url: str
    mermaid_svg_url: str
    validation_feedback: List[str]
    ai_suggestions: Dict[str, Any]


def get_similar_program_patterns(theme: str):
    return ecosystem_patterns_collection.find_one(
        {"theme": theme},
        {"_id": 0}
    )

def get_district_challenges(state: str, district: str):
    record = district_challenges_collection.find_one(
        {"state": state.lower(), "district": district.lower()},
        {"_id": 0}
    )
    return record["challenges"] if record else []

def validate_problem_against_ecosystem(problem_text: str, pattern: Dict[str, Any]):
    feedback = []

    if pattern:
        if pattern["core_problem_pattern"].lower() not in problem_text.lower():
            feedback.append(
                "Your problem statement deviates from common ecosystem patterns. Ensure this is intentional."
            )

    if len(problem_text.split()) < 10:
        feedback.append("Problem statement may be too short for systemic clarity.")

    return feedback

def build_problem_tree_structure(
    refined_problem_statement: str,
    root_causes: List[Dict[str, str]],
    ecosystem_pattern: Optional[Dict[str, Any]]
):
    causes = []
    effects = []

    for i, cause in enumerate(root_causes):
        causes.append({
            "id": f"C{i+1}",
            "label": cause["cause"]
        })

    if ecosystem_pattern:
        for ec in ecosystem_pattern.get("common_effects", []):
            effects.append({
                "id": f"E{len(effects)+1}",
                "label": ec
            })

    effects.append({
        "id": f"E{len(effects)+1}",
        "label": "Limited long-term systemic impact"
    })

    return {
        "causes": causes,
        "core_problem": {
            "id": "P1",
            "label": refined_problem_statement
        },
        "effects": effects
    }

def generate_mermaid_problem_tree(tree: Dict[str, Any]) -> str:
    lines = ["graph LR"]

    for cause in tree["causes"]:
        lines.append(f'{cause["id"]}["{cause["label"]}"] --> P1')

    lines.append(f'P1["{tree["core_problem"]["label"]}"]')

    for effect in tree["effects"]:
        lines.append(f'P1 --> {effect["id"]}["{effect["label"]}"]')

    return "\n".join(lines)

def mermaid_to_live_url(mermaid_code: str) -> str:
    payload = {
        "code": mermaid_code,
        "mermaid": {"theme": "default"}
    }
    encoded = base64.urlsafe_b64encode(
        json.dumps(payload).encode("utf-8")
    ).decode("utf-8")

    return f"https://mermaid.live/edit#{encoded}"

def mermaid_to_image_url(mermaid_code: str, format: str = "png") -> str:
    """
    format: 'png' or 'svg'
    """
    encoded = base64.urlsafe_b64encode(
        mermaid_code.encode("utf-8")
    ).decode("utf-8")

    return f"https://mermaid.ink/img/{encoded}?type={format}"

#  PROBLEM TREE GENERATION API 
@app.post("/problem-tree", response_model=ProblemTreeResponse)
def generate_problem_tree(payload: ProblemTreeRequest):

    ecosystem_pattern = get_similar_program_patterns(payload.theme)
    district_challenges = get_district_challenges(payload.state, payload.district)

    problem_tree = build_problem_tree_structure(
        payload.refined_problem_statement,
        payload.suggested_root_causes,
        ecosystem_pattern
    )

    mermaid_diagram = generate_mermaid_problem_tree(problem_tree)
    mermaid_preview_url = mermaid_to_live_url(mermaid_diagram)
    mermaid_png_url = mermaid_to_image_url(mermaid_diagram, "png")
    mermaid_svg_url = mermaid_to_image_url(mermaid_diagram, "svg")

    validation_feedback = validate_problem_against_ecosystem(
        payload.refined_problem_statement,
        ecosystem_pattern
    )

    ai_suggestions = {
        "similar_program_pattern": ecosystem_pattern,
        "district_challenges": district_challenges
    }

    record = {
        "id": str(uuid.uuid4()),
        "organization_id": payload.organization_id,
        "problem_tree": problem_tree,
        "theme": payload.theme,
        "state": payload.state,
        "district": payload.district,
        "created_at": datetime.utcnow()
    }

    problem_tree_collection.insert_one(record)

    return {
    "problem_tree": problem_tree,
    "mermaid_diagram": mermaid_diagram,
    "mermaid_preview_url": mermaid_preview_url,
    "mermaid_png_url": mermaid_png_url,
    "mermaid_svg_url": mermaid_svg_url,
    "validation_feedback": validation_feedback,
    "ai_suggestions": ai_suggestions
}

# -------------------- Student Outcome Designer --------------------
class StudentOutcomeRequest(BaseModel):
    organization_id: str
    theme: str
    state: str
    grade_range: str
    outcome_statement: str
    baseline_value: float
    target_value: float
    timeline_months: int

class StudentOutcomeResponse(BaseModel):
    smart_validation: Dict[str, Any]
    aligned_competencies: List[str]
    policy_references: List[str]
    outcome_timeline_diagram: str
    timeline_preview_url: str
    timeline_png_url: str
    timeline_svg_url: str

def validate_smart_outcome(
    outcome: str,
    baseline: float,
    target: float,
    timeline: int
):

    feedback = []
    score = 5

    if len(outcome.split()) < 8:
        feedback.append("Outcome may not be sufficiently specific.")
        score -= 1

    if baseline >= target:
        feedback.append("Target value must be greater than baseline.")
        score -= 1

    if timeline <= 0:
        feedback.append("Timeline must be greater than zero months.")
        score -= 1

    measurable = any(word in outcome.lower() for word in ["percentage", "score", "proficiency", "fluency"])

    if not measurable:
        feedback.append("Outcome may not be clearly measurable.")
        score -= 1

    return {
        "smart_score": max(score, 1),
        "issues": feedback,
        "is_valid": score >= 3
    }

def get_aligned_competencies(theme: str, grade_range: str):
    record = competency_frameworks_collection.find_one(
        {"theme": theme, "grade_range": grade_range},
        {"_id": 0}
    )
    return record["competencies"] if record else []

def get_policy_references():
    return [
        p["reference"]
        for p in policy_references_collection.find({}, {"_id": 0})
    ]

def generate_outcome_timeline(outcome: str, timeline_months: int) -> str:

    return f"""
gantt
    title Student Outcome Timeline
    dateFormat  YYYY-MM-DD
    section Outcome Achievement
    Baseline to Target Progress :a1, 2025-01-01, {timeline_months * 30}d
"""
def mermaid_to_live_url(mermaid_code: str) -> str:
    payload = {
        "code": mermaid_code,
        "mermaid": {"theme": "default"}
    }
    encoded = base64.urlsafe_b64encode(
        json.dumps(payload).encode("utf-8")
    ).decode("utf-8")

    return f"https://mermaid.live/edit#{encoded}"

def mermaid_to_image_url(mermaid_code: str, format: str = "png") -> str:
    encoded = base64.urlsafe_b64encode(
        mermaid_code.encode("utf-8")
    ).decode("utf-8")

    return f"https://mermaid.ink/img/{encoded}?type={format}"

# STUDENT OUTCOME CREATION API
@app.post("/student-outcomes", response_model=StudentOutcomeResponse)
def create_student_outcome(payload: StudentOutcomeRequest):

    smart_validation = validate_smart_outcome(
        payload.outcome_statement,
        payload.baseline_value,
        payload.target_value,
        payload.timeline_months
    )

    aligned_competencies = get_aligned_competencies(
        payload.theme,
        payload.grade_range
    )

    policy_refs = get_policy_references()

    timeline_diagram = generate_outcome_timeline(
        payload.outcome_statement,
        payload.timeline_months
    )

    # 👇 Visualization support
    timeline_preview_url = mermaid_to_live_url(timeline_diagram)
    timeline_png_url = mermaid_to_image_url(timeline_diagram, "png")
    timeline_svg_url = mermaid_to_image_url(timeline_diagram, "svg")

    record = {
        "id": str(uuid.uuid4()),
        "organization_id": payload.organization_id,
        "theme": payload.theme,
        "state": payload.state,
        "grade_range": payload.grade_range,
        "outcome_statement": payload.outcome_statement,
        "baseline": payload.baseline_value,
        "target": payload.target_value,
        "timeline_months": payload.timeline_months,
        "smart_validation": smart_validation,
        "competencies": aligned_competencies,
        "policy_references": policy_refs,
        "created_at": datetime.now(timezone.utc)
    }

    student_outcomes_collection.insert_one(record)

    return {
        "smart_validation": smart_validation,
        "aligned_competencies": aligned_competencies,
        "policy_references": policy_refs,
        "outcome_timeline_diagram": timeline_diagram,
        "timeline_preview_url": timeline_preview_url,
        "timeline_png_url": timeline_png_url,
        "timeline_svg_url": timeline_svg_url
    }

# -------------------- Methodology Selector --------------------
class MethodologyFilterRequest(BaseModel):
    theme: str
    state: str
    scale_schools: int
    budget_lakhs: int

class MethodologyResponse(BaseModel):
    methodologies: List[Dict[str, Any]]

def filter_methodologies(
    theme: str,
    state: str,
    scale: int,
    budget: int
):

    results = []

    cursor = methodology_library_collection.find(
        {"theme": theme},
        {"_id": 0}
    )

    for m in cursor:
        geo_match = state.lower() in m["geographies"] or "all" in m["geographies"]

        min_budget, max_budget = m["budget_range_lakhs"]
        budget_match = min_budget <= budget <= max_budget

        if geo_match and budget_match:
            results.append(m)

    return results

def generate_component_library(methodologies: List[Dict[str, Any]]):

    component_map = {}

    for m in methodologies:
        for component in m["components"]:
            component_map.setdefault(component, []).append(m["name"])

    return [
        {
            "component": c,
            "used_in": sources
        }
        for c, sources in component_map.items()
    ]

# METHODOLOGY SELECTION API
@app.post("/methodologies", response_model=MethodologyResponse)
def get_methodologies(payload: MethodologyFilterRequest):

    methodologies = filter_methodologies(
        payload.theme,
        payload.state,
        payload.scale_schools,
        payload.budget_lakhs
    )

    component_library = generate_component_library(methodologies)

    return {
        "methodologies": [
            {
                **m,
                "available_components": component_library
            }
            for m in methodologies
        ]
    }

class SelectMethodologyRequest(BaseModel):
    organization_id: str
    selected_methodology_ids: List[str]
    custom_components: List[str]

# METHODOLOGY SAVE SELECTION API
@app.post("/methodologies/select")
def save_selected_methodology(payload: SelectMethodologyRequest):

    record = {
        "organization_id": payload.organization_id,
        "selected_methodology_ids": payload.selected_methodology_ids,
        "custom_components": payload.custom_components,
        "created_at": datetime.utcnow()
    }

    selected_methodologies_collection.insert_one(record)

    return {"status": "Methodologies saved successfully"}

# -------------------- THEORY OF CHANGE BUILDER --------------------
class ToCNode(BaseModel):
    id: str
    type: str  # activity | output | outcome | impact
    label: str

class ToCEdge(BaseModel):
    source: str
    target: str

class TheoryOfChangeRequest(BaseModel):
    organization_id: str
    theme: str
    nodes: List[ToCNode]
    edges: List[ToCEdge]

class TheoryOfChangeResponse(BaseModel):
    is_valid: bool
    logic_issues: List[str]
    ai_suggestions: List[str]
    mermaid_diagram: str
    mermaid_preview_url: str
    mermaid_png_url: str
    mermaid_svg_url: str

def validate_if_then_logic(nodes: List[ToCNode], edges: List[ToCEdge]):

    issues = []
    node_map = {node.id: node for node in nodes}

    required_flow = ["activity", "output", "outcome", "impact"]

    # Map outgoing edges
    flow_map = {}
    for edge in edges:
        flow_map.setdefault(edge.source, []).append(edge.target)

    # Validate each connection
    for edge in edges:
        source_type = node_map[edge.source].type
        target_type = node_map[edge.target].type

        if required_flow.index(target_type) - required_flow.index(source_type) != 1:
            issues.append(
                f"Invalid logical jump from {source_type} to {target_type}: '{node_map[edge.source].label}' → '{node_map[edge.target].label}'"
            )

    # Check missing layers
    present_layers = {node.type for node in nodes}
    for layer in required_flow:
        if layer not in present_layers:
            issues.append(f"Missing '{layer}' level in Theory of Change.")

    return issues

def detect_logic_gaps(theme: str, nodes: List[ToCNode]):

    suggestions = []
    pattern = toc_pattern_library_collection.find_one({"theme": theme})

    if not pattern:
        return ["No reference ToC pattern found for this theme."]

    node_labels = [n.label.lower() for n in nodes]

    for expected_output in pattern["outputs"]:
        if expected_output.lower() not in node_labels:
            suggestions.append(
                f"Consider adding output: '{expected_output}' (commonly seen in successful programs)."
            )

    for expected_outcome in pattern["outcomes"]:
        if expected_outcome.lower() not in node_labels:
            suggestions.append(
                f"Missing typical outcome: '{expected_outcome}'."
            )

    return suggestions

def generate_toc_mermaid(nodes: List[ToCNode], edges: List[ToCEdge]):

    lines = ["flowchart LR"]

    for node in nodes:
        label = (
            node.label
            .replace('"', "'")
            .replace("[", "")
            .replace("]", "")
            .replace("\n", " ")
        )

        if node.type == "activity":
            lines.append(f'{node.id}["{label}"]')
        elif node.type == "output":
            lines.append(f'{node.id}("{label}")')
        elif node.type == "outcome":
            lines.append(f'{node.id}(("{label}"))')
        elif node.type == "impact":
            lines.append(f'{node.id}((("{label}")))')
        else:
            lines.append(f'{node.id}["{label}"]')

    for edge in edges:
        lines.append(f"{edge.source} --> {edge.target}")

    return "\n".join(lines)

# THEORY OF CHANGE BUILDING API
@app.post("/theory-of-change", response_model=TheoryOfChangeResponse)
def build_theory_of_change(payload: TheoryOfChangeRequest):

    logic_issues = validate_if_then_logic(payload.nodes, payload.edges)
    ai_suggestions = detect_logic_gaps(payload.theme, payload.nodes)

    is_valid = len(logic_issues) == 0

    mermaid_diagram = generate_toc_mermaid(
        payload.nodes,
        payload.edges
    )

    # 👇 Visualization support
    mermaid_preview_url = mermaid_to_live_url(mermaid_diagram)
    mermaid_png_url = mermaid_to_image_url(mermaid_diagram, "png")
    mermaid_svg_url = mermaid_to_image_url(mermaid_diagram, "svg")

    record = {
        "organization_id": payload.organization_id,
        "theme": payload.theme,
        "nodes": [n.dict() for n in payload.nodes],
        "edges": [e.dict() for e in payload.edges],
        "logic_issues": logic_issues,
        "ai_suggestions": ai_suggestions,
        "created_at": datetime.utcnow()
    }

    theory_of_change_collection.insert_one(record)

    return {
        "is_valid": is_valid,
        "logic_issues": logic_issues,
        "ai_suggestions": ai_suggestions,
        "mermaid_diagram": mermaid_diagram,
        "mermaid_preview_url": mermaid_preview_url,
        "mermaid_png_url": mermaid_png_url,
        "mermaid_svg_url": mermaid_svg_url
    }

# -------------------- Stakeholder Selector --------------------
class StakeholderSelectorRequest(BaseModel):
    organization_id: str
    theme: str

class StakeholderSelectorResponse(BaseModel):
    available_stakeholders: List[Dict[str, Any]]
    recommended_stakeholders: List[str]

def get_recommended_stakeholders(theme: str):

    cursor = stakeholder_master_collection.find(
        {"themes": theme},
        {"_id": 0}
    )

    return [s["stakeholder_id"] for s in cursor]

# STAKEHOLDER SELECTION API
@app.post("/stakeholders/select", response_model=StakeholderSelectorResponse)
def select_stakeholders(payload: StakeholderSelectorRequest):

    all_stakeholders = list(
        stakeholder_master_collection.find({}, {"_id": 0})
    )

    recommended = get_recommended_stakeholders(payload.theme)

    record = {
        "organization_id": payload.organization_id,
        "theme": payload.theme,
        "selected_stakeholders": recommended,
        "created_at": datetime.utcnow()
    }

    organization_stakeholders_collection.insert_one(record)

    return {
        "available_stakeholders": all_stakeholders,
        "recommended_stakeholders": recommended
    }

# --------------------  Practice Change Definition Tool --------------------
class PracticeChangeRequest(BaseModel):
    organization_id: str
    theme: str
    stakeholder_id: str
    current_practices: List[str]
    desired_practices: List[str]

class PracticeChangeResponse(BaseModel):
    ai_suggestions: Dict[str, List[str]]
    validation_feedback: List[str]

def get_ai_practice_suggestions(stakeholder_id: str, theme: str):

    record = practice_master_collection.find_one(
        {"stakeholder_id": stakeholder_id, "theme": theme},
        {"_id": 0}
    )

    if not record:
        return {"suggested_current": [], "suggested_desired": []}

    return {
        "suggested_current": record.get("current_practices", []),
        "suggested_desired": record.get("desired_practices", [])
    }

def validate_practice_changes(current: List[str], desired: List[str]):

    feedback = []

    if not current:
        feedback.append("Current practices cannot be empty.")

    if not desired:
        feedback.append("Desired practices cannot be empty.")

    if len(desired) < len(current):
        feedback.append(
            "Desired practices should demonstrate improvement beyond current practices."
        )

    vague_terms = ["better", "improved", "enhanced", "effective"]

    for practice in desired:
        if any(v in practice.lower() for v in vague_terms):
            feedback.append(
                f"Desired practice '{practice}' may be too vague. Consider making it more specific."
            )

    return feedback

# PRACTICE CHANGE DEFINITION API
@app.post("/practice-change", response_model=PracticeChangeResponse)
def define_practice_change(payload: PracticeChangeRequest):

    ai_suggestions = get_ai_practice_suggestions(
        payload.stakeholder_id,
        payload.theme
    )

    validation_feedback = validate_practice_changes(
        payload.current_practices,
        payload.desired_practices
    )

    record = {
        "organization_id": payload.organization_id,
        "theme": payload.theme,
        "stakeholder_id": payload.stakeholder_id,
        "current_practices": payload.current_practices,
        "desired_practices": payload.desired_practices,
        "created_at": datetime.utcnow()
    }

    practice_change_collection.insert_one(record)

    return {
        "ai_suggestions": ai_suggestions,
        "validation_feedback": validation_feedback
    }

# -------------------- Auto Indicator Suggestion Engine  --------------------
class IndicatorGenerationRequest(BaseModel):
    organization_id: str
    theme: str
    student_outcomes: List[str]
    practice_changes: List[Dict[str, Any]]
    # practice_changes format:
    # {
    #   stakeholder_id: "TCH",
    #   desired_practices: [ "...", "..." ]
    # }

class IndicatorGenerationResponse(BaseModel):
    outcome_indicators: Dict[str, str]
    practice_indicators: Dict[str, Dict[str, str]]

def generate_outcome_indicators(theme: str, outcomes: List[str]):

    record = indicator_master_collection.find_one(
        {"type": "student_outcome", "theme": theme},
        {"_id": 0}
    )

    templates = record["indicator_templates"] if record else []

    indicators = {}

    for i, outcome in enumerate(outcomes):
        if templates:
            indicators[outcome] = templates[i % len(templates)]
        else:
            indicators[outcome] = f"% of students achieving '{outcome}'"

    return indicators

def generate_practice_indicators(theme: str, practice_changes: List[Dict[str, Any]]):

    practice_indicators = {}

    for pc in practice_changes:
        stakeholder = pc["stakeholder_id"]
        desired_practices = pc["desired_practices"]

        record = indicator_master_collection.find_one(
            {
                "type": "practice_change",
                "stakeholder_id": stakeholder,
                "theme": theme
            },
            {"_id": 0}
        )

        templates = record["indicator_templates"] if record else []

        stakeholder_indicators = {}

        for i, practice in enumerate(desired_practices):
            if templates:
                stakeholder_indicators[practice] = templates[i % len(templates)]
            else:
                stakeholder_indicators[practice] = f"% adoption of practice: '{practice}'"

        practice_indicators[stakeholder] = stakeholder_indicators

    return practice_indicators

# AUTO INDICATOR GENERATION API
@app.post("/indicators/auto-generate", response_model=IndicatorGenerationResponse)
def auto_generate_indicators(payload: IndicatorGenerationRequest):

    outcome_indicators = generate_outcome_indicators(
        payload.theme,
        payload.student_outcomes
    )

    practice_indicators = generate_practice_indicators(
        payload.theme,
        payload.practice_changes
    )

    record = {
        "organization_id": payload.organization_id,
        "theme": payload.theme,
        "outcome_indicators": outcome_indicators,
        "practice_indicators": practice_indicators,
        "created_at": datetime.utcnow()
    }

    generated_indicators_collection.insert_one(record)

    return {
        "outcome_indicators": outcome_indicators,
        "practice_indicators": practice_indicators
    }

# --------------------  Baseline, Target & Timeline Builder--------------------
class IndicatorTarget(BaseModel):
    indicator_name: str
    baseline_value: float
    target_value: float
    start_date: datetime
    end_date: datetime

class BaselineTargetRequest(BaseModel):
    organization_id: str
    indicators: List[IndicatorTarget]

class ValidationResult(BaseModel):
    indicator_name: str
    status: str
    warnings: List[str]

class BaselineTargetResponse(BaseModel):
    validations: List[ValidationResult]

def validate_indicator_target(indicator: IndicatorTarget):

    warnings = []

    # Rule 1: Target must be greater than baseline
    if indicator.target_value <= indicator.baseline_value:
        warnings.append("Target value must be greater than baseline")

    # Rule 2: Timeline check
    duration_months = (
        (indicator.end_date.year - indicator.start_date.year) * 12 +
        (indicator.end_date.month - indicator.start_date.month)
    )

    if duration_months <= 0:
        warnings.append("End date must be after start date")

    # Rule 3: Unrealistic improvement check
    improvement = indicator.target_value - indicator.baseline_value

    if duration_months > 0:
        monthly_growth = improvement / duration_months

        if monthly_growth > 10:
            warnings.append(
                "Target growth appears aggressive for the given timeline"
            )

    status = "valid" if not warnings else "needs_review"

    return ValidationResult(
        indicator_name=indicator.indicator_name,
        status=status,
        warnings=warnings
    )

# BASELINE & TARGET SETTING API
@app.post(
    "/indicators/baseline-target",
    response_model=BaselineTargetResponse
)
def build_baseline_target(payload: BaselineTargetRequest):

    validations = []

    for indicator in payload.indicators:

        validation = validate_indicator_target(indicator)
        validations.append(validation)

        record = {
            "organization_id": payload.organization_id,
            "indicator_name": indicator.indicator_name,
            "baseline_value": indicator.baseline_value,
            "target_value": indicator.target_value,
            "start_date": indicator.start_date,
            "end_date": indicator.end_date,
            "status": validation.status,
            "warnings": validation.warnings,
            "created_at": datetime.utcnow()
        }

        indicator_targets_collection.insert_one(record)

    return {"validations": validations}

# --------------------  AI LFA COMPLETENESS SCORE ENGINE --------------------
SECTION_REQUIREMENTS = {
    "organization_profile": {
        "weight": 10,
        "required_fields": ["organization_id", "theme", "geography", "scale"]
    },
    "problem_definition": {
        "weight": 15,
        "required_fields": ["core_problem", "affected_group"]
    },
    "problem_tree": {
        "weight": 15,
        "required_fields": ["root_causes", "core_problem", "effects"]
    },
    "outcomes": {
        "weight": 20,
        "required_fields": ["smart_outcomes"]
    },
    "methodology": {
        "weight": 15,
        "required_fields": ["interventions"]
    },
    "theory_of_change": {
        "weight": 15,
        "required_fields": ["activities", "outputs", "outcomes", "impact"]
    },
    "measurement": {
        "weight": 10,
        "required_fields": ["indicators", "targets"]
    }
}

class LFACompletenessRequest(BaseModel):
    organization_id: str
    lfa_snapshot: Dict[str, Any]

def calculate_lfa_completeness(lfa_snapshot: Dict[str, Any]):

    section_results = []
    total_score = 0
    missing_sections = []

    for section, rules in SECTION_REQUIREMENTS.items():
        section_data = lfa_snapshot.get(section)

        if not section_data:
            missing_sections.append(section)
            section_results.append({
                "section": section,
                "status": "missing",
                "score": 0,
                "weight": rules["weight"]
            })
            continue

        # Check required fields
        present_fields = [
            field for field in rules["required_fields"]
            if field in section_data and section_data[field]
        ]

        completion_ratio = len(present_fields) / len(rules["required_fields"])
        section_score = round(rules["weight"] * completion_ratio, 2)

        total_score += section_score

        status = "complete" if completion_ratio == 1 else "partial"

        section_results.append({
            "section": section,
            "status": status,
            "score": section_score,
            "weight": rules["weight"],
            "missing_fields": list(
                set(rules["required_fields"]) - set(present_fields)
            )
        })

    completion_percentage = round(total_score, 2)

    return {
        "completion_percentage": completion_percentage,
        "section_breakdown": section_results,
        "missing_sections": missing_sections
    }

# LFA COMPLETENESS SCORE API
@app.post("/lfa/completeness-score")
def get_lfa_completeness(payload: LFACompletenessRequest):

    result = calculate_lfa_completeness(payload.lfa_snapshot)

    record = {
        "organization_id": payload.organization_id,
        "completion_percentage": result["completion_percentage"],
        "section_breakdown": result["section_breakdown"],
        "missing_sections": result["missing_sections"],
        "evaluated_at": datetime.utcnow()
    }

    lfa_completeness_collection.insert_one(record)

    return result

# -------------------- DESIGN QUALITY FEEDBACK SYSTEM --------------------
QUALITY_CHECKS = {
    "outcome_quality": {
        "severity": "high"
    },
    "stakeholder_alignment": {
        "severity": "high"
    },
    "indicator_validity": {
        "severity": "medium"
    },
    "theory_of_change_logic": {
        "severity": "high"
    },
    "problem_intervention_alignment": {
        "severity": "medium"
    }
}

class DesignQualityRequest(BaseModel):
    organization_id: str
    lfa_snapshot: Dict[str, Any]

def analyze_outcome_quality(outcomes: List[str]):

    feedback = []

    measurable_keywords = ["increase", "decrease", "%", "to", "by", "from"]
    vague_words = ["improve", "enhance", "strengthen", "better"]

    for outcome in outcomes:
        issues = []

        if not any(word in outcome.lower() for word in measurable_keywords):
            issues.append("Outcome is not clearly measurable")

        if any(word in outcome.lower() for word in vague_words):
            issues.append("Outcome uses vague language")

        if issues:
            feedback.append({
                "area": "Student Outcomes",
                "issue": outcome,
                "problems": issues,
                "suggestion": "Rewrite outcome using a measurable baseline, target, and timeframe"
            })

    return feedback

def analyze_stakeholder_alignment(stakeholders, outcomes):

    feedback = []

    if not stakeholders:
        feedback.append({
            "area": "Stakeholders",
            "issue": "No stakeholders linked to outcomes",
            "suggestion": "Map at least one stakeholder responsible for each outcome"
        })
        return feedback

    if outcomes and len(stakeholders) < len(outcomes):
        feedback.append({
            "area": "Stakeholder Coverage",
            "issue": "Fewer stakeholders than outcomes",
            "suggestion": "Ensure accountability by mapping stakeholders to each outcome"
        })

    return feedback

def analyze_indicator_quality(indicators):

    feedback = []

    for indicator in indicators:
        if "survey" in indicator.lower() or "perception" in indicator.lower():
            feedback.append({
                "area": "Indicators",
                "issue": indicator,
                "problem": "Indicator is perception-based",
                "suggestion": "Prefer observable or performance-based indicators"
            })

    return feedback

def analyze_toc_logic(toc):

    feedback = []

    required_chain = ["activities", "outputs", "outcomes", "impact"]

    for step in required_chain:
        if step not in toc or not toc[step]:
            feedback.append({
                "area": "Theory of Change",
                "issue": f"Missing {step}",
                "suggestion": f"Define clear {step} to maintain logical flow"
            })

    if "activities" in toc and "outcomes" in toc:
        if len(toc["activities"]) > len(toc["outcomes"]) * 3:
            feedback.append({
                "area": "Theory of Change",
                "issue": "Too many activities for defined outcomes",
                "suggestion": "Reduce activities or clarify outcome pathways"
            })

    return feedback

def analyze_problem_intervention_alignment(problem, interventions):

    feedback = []

    if problem and interventions:
        mismatch_terms = ["infrastructure", "hardware"]

        if any(term in problem.lower() for term in mismatch_terms):
            if not any("infrastructure" in i.lower() for i in interventions):
                feedback.append({
                    "area": "Intervention Alignment",
                    "issue": "Problem–intervention mismatch",
                    "suggestion": "Selected interventions do not address the stated core problem"
                })

    return feedback

def generate_design_quality_feedback(lfa_snapshot):

    feedback = []

    outcomes = lfa_snapshot.get("outcomes", {}).get("smart_outcomes", [])
    stakeholders = lfa_snapshot.get("stakeholders", [])
    indicators = lfa_snapshot.get("measurement", {}).get("indicators", [])
    toc = lfa_snapshot.get("theory_of_change", {})
    problem = lfa_snapshot.get("problem_definition", {}).get("core_problem", "")
    interventions = lfa_snapshot.get("methodology", {}).get("interventions", [])

    feedback.extend(analyze_outcome_quality(outcomes))
    feedback.extend(analyze_stakeholder_alignment(stakeholders, outcomes))
    feedback.extend(analyze_indicator_quality(indicators))
    feedback.extend(analyze_toc_logic(toc))
    feedback.extend(analyze_problem_intervention_alignment(problem, interventions))

    quality_score = max(0, 100 - len(feedback) * 8)

    return {
        "quality_score": quality_score,
        "feedback_items": feedback
    }

# DESIGN QUALITY FEEDBACK API
@app.post("/lfa/design-quality-feedback")
def get_design_quality_feedback(payload: DesignQualityRequest):

    result = generate_design_quality_feedback(payload.lfa_snapshot)

    record = {
        "organization_id": payload.organization_id,
        "quality_score": result["quality_score"],
        "feedback_items": result["feedback_items"],
        "evaluated_at": datetime.utcnow()
    }

    design_quality_feedback_collection.insert_one(record)

    return result

# -------------------- COMMON LFA TEMPLATE ENGINE --------------------
class LFATemplate(BaseModel):
    template_id: str
    name: str
    theme: str
    system_level: str
    geography_type: str  # state / urban / rural / tribal / generic
    applicable_states: Optional[List[str]] = []
    description: str
    lfa_structure: Dict[str, Any]
    created_by: str  # "system" or organization_id
    is_public: bool
    created_at: datetime

# LFA TEMPLATE LISTING API
@app.get("/lfa/templates")
def list_lfa_templates(
    theme: Optional[str] = None,
    system_level: Optional[str] = None,
    geography_type: Optional[str] = None
):

    query = {"is_public": True}

    if theme:
        query["theme"] = theme
    if system_level:
        query["system_level"] = system_level
    if geography_type:
        query["geography_type"] = geography_type

    templates = list(
        lfa_templates_collection.find(query, {"_id": 0})
    )

    return {
        "count": len(templates),
        "templates": templates
    }

class ForkTemplateRequest(BaseModel):
    organization_id: str
    template_id: str
    new_name: str

# LFA TEMPLATE FORKING API
@app.post("/lfa/templates/fork")
def fork_lfa_template(payload: ForkTemplateRequest):

    base_template = lfa_templates_collection.find_one(
        {"template_id": payload.template_id},
        {"_id": 0}
    )

    if not base_template:
        raise HTTPException(status_code=404, detail="Template not found")

    forked_template = base_template.copy()
    forked_template["template_id"] = f"{payload.organization_id}_{int(datetime.now(timezone.utc).timestamp())}"
    forked_template["name"] = payload.new_name
    forked_template["created_by"] = payload.organization_id
    forked_template["is_public"] = False
    forked_template["created_at"] = datetime.utcnow()

    organization_templates_collection.insert_one(forked_template)

    return {
        "message": "Template forked successfully",
        "template_id": forked_template["template_id"]
    }

# LFA CUSTOM TEMPLATE SAVE API
@app.post("/lfa/templates/save")
def save_custom_template(template: LFATemplate):

    organization_templates_collection.insert_one(template.dict())

    return {
        "message": "Custom template saved successfully",
        "template_id": template.template_id
    }

# LFA TEMPLATE SHARE TO MARKETPLACE API
@app.post("/lfa/templates/share/{template_id}")
def share_template(template_id: str):

    result = organization_templates_collection.update_one(
        {"template_id": template_id},
        {"$set": {"is_public": True}}
    )

    if result.modified_count == 0:
        raise HTTPException(status_code=404, detail="Template not found")

    return {"message": "Template shared to marketplace"}

class TemplateRatingRequest(BaseModel):
    template_id: str
    organization_id: str
    rating: int = Field(..., ge=1, le=5)

# LFA TEMPLATE RATING API
@app.post("/lfa/templates/rate")
def rate_template(payload: TemplateRatingRequest):

    template_ratings_collection.insert_one({
        "template_id": payload.template_id,
        "organization_id": payload.organization_id,
        "rating": payload.rating,
        "rated_at": datetime.utcnow()
    })

    return {"message": "Rating submitted"}

# --------------------  EXPORT & INTEGRATION ENGINE --------------------
EXPORT_TYPES = [
    "LFA_PDF",
    "TOC_IMAGE",
    "INDICATOR_EXCEL",
    "BUDGET_EXCEL",
    "PRESENTATION_PPT"
]

class ExportRequest(BaseModel):
    organization_id: str
    export_type: str
    lfa_snapshot: Dict[str, Any]

def generate_lfa_pdf(lfa_snapshot):

    file_path = f"/tmp/LFA_{int(datetime.now(timezone.utc).timestamp())}.pdf"
    doc = SimpleDocTemplate(file_path)

    styles = getSampleStyleSheet()
    content = []

    content.append(Paragraph("Logical Framework Matrix", styles["Title"]))
    content.append(Spacer(1, 12))

    for section, data in lfa_snapshot.items():
        content.append(Paragraph(section.replace("_", " ").title(), styles["Heading2"]))
        content.append(Paragraph(json.dumps(data, indent=2), styles["Normal"]))
        content.append(Spacer(1, 10))

    doc.build(content)
    return file_path

def generate_toc_diagram(toc: dict):

    G = nx.DiGraph()

    activities = toc.get("activities", [])
    outputs = toc.get("outputs", [])
    outcomes = toc.get("outcomes", [])
    impacts = toc.get("impact", [])

    # --- Add nodes with layers ---
    for a in activities:
        G.add_node(a, layer=0)

    for o in outputs:
        G.add_node(o, layer=1)

    for o in outcomes:
        G.add_node(o, layer=2)

    for i in impacts:
        G.add_node(i, layer=3)

    # --- Add edges (logical flow) ---
    for a in activities:
        for o in outputs:
            G.add_edge(a, o)

    for o in outputs:
        for oc in outcomes:
            G.add_edge(o, oc)

    for oc in outcomes:
        for i in impacts:
            G.add_edge(oc, i)

    # --- Layout: left → right by layer ---
    pos = {}
    layer_nodes = {}

    for node, data in G.nodes(data=True):
        layer = data["layer"]
        layer_nodes.setdefault(layer, []).append(node)

    for layer, nodes in layer_nodes.items():
        for i, node in enumerate(nodes):
            pos[node] = (layer * 4, -i)

    # --- Draw ---
    plt.figure(figsize=(14, 8))
    nx.draw(
        G,
        pos,
        with_labels=True,
        node_size=3000,
        node_shape="s",
        font_size=9,
        font_weight="bold"
    )

    file_path = f"/tmp/TOC_{int(time.time())}.png"
    plt.tight_layout()
    plt.savefig(file_path)
    plt.close()

    return file_path

def generate_indicator_excel(measurement):

    file_path = f"/tmp/Indicators_{int(datetime.now(timezone.utc).timestamp())}.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Indicators"

    ws.append(["Outcome", "Indicator", "Baseline", "Target", "Timeline"])

    for item in measurement.get("indicators", []):
        ws.append([
            item.get("outcome"),
            item.get("indicator"),
            item.get("baseline"),
            item.get("target"),
            item.get("timeline")
        ])

    wb.save(file_path)
    return file_path

def generate_budget_template():

    file_path = f"/tmp/Budget_{int(datetime.now(timezone.utc).timestamp())}.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Budget"

    ws.append(["Component", "Unit Cost", "Quantity", "Total", "Notes"])

    wb.save(file_path)
    return file_path

def generate_presentation(lfa_snapshot):

    file_path = f"/tmp/Program_Overview_{int(datetime.now(timezone.utc).timestamp())}.pptx"
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Program Design Overview"

    for section, data in lfa_snapshot.items():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section.replace("_", " ").title()
        slide.placeholders[1].text = json.dumps(data, indent=2)

    prs.save(file_path)
    return file_path

def handle_export(export_type, lfa_snapshot):

    if export_type == "LFA_PDF":
        return generate_lfa_pdf(lfa_snapshot)

    if export_type == "TOC_IMAGE":
        return generate_toc_diagram(
            lfa_snapshot.get("theory_of_change", {})
        )

    if export_type == "INDICATOR_EXCEL":
        return generate_indicator_excel(
            lfa_snapshot.get("measurement", {})
        )

    if export_type == "BUDGET_EXCEL":
        return generate_budget_template()

    if export_type == "PRESENTATION_PPT":
        return generate_presentation(lfa_snapshot)

    raise HTTPException(status_code=400, detail="Unsupported export type")

# EXPORT API
@app.post("/export")
def export_lfa(payload: ExportRequest):

    if payload.export_type not in EXPORT_TYPES:
        raise HTTPException(status_code=400, detail="Invalid export type")

    file_path = handle_export(payload.export_type, payload.lfa_snapshot)

    export_jobs_collection.insert_one({
        "organization_id": payload.organization_id,
        "export_type": payload.export_type,
        "file_path": file_path,
        "created_at": datetime.utcnow()
    })

    return FileResponse(
        file_path,
        filename=file_path.split("/")[-1],
        media_type="application/octet-stream"
    )

# -------------------- Multi Lingual Support--------------------
translator = Translator()

def translate_object(obj: Any, target_lang: str = "en") -> Any:
    """
    Recursively translate all strings in dict, list, or str.
    """
    if isinstance(obj, str):
        if target_lang == "en":
            return obj
        try:
            return translator.translate(obj, dest=target_lang).text
        except Exception:
            return obj
    elif isinstance(obj, dict):
        return {k: translate_object(v, target_lang) for k, v in obj.items()}
    elif isinstance(obj, list):
        return [translate_object(i, target_lang) for i in obj]
    else:
        return obj

def translate_response(response_obj: Any, language: str = "en") -> Any:
    """
    Translate all strings in the response object to the target language
    """
    if language == "en":
        return response_obj
    return translate_object(response_obj, language)

# Translation Middleware 
@app.middleware("http")
async def translate_middleware(request: Request, call_next):
    # 1️⃣ Capture the 'language' query parameter (default: English)
    language = request.query_params.get("language", "en")
    
    # 2️⃣ Process the request normally
    response = await call_next(request)
    
    # 3️⃣ Only translate JSON responses
    if "application/json" in response.headers.get("content-type", ""):
        body = [section async for section in response.body_iterator]
        if body:
            raw = b"".join(body).decode()
            try:
                data = json.loads(raw)
            except json.JSONDecodeError:
                return response
            translated = translate_response(data, language)
            return JSONResponse(content=translated, status_code=response.status_code)
    
    # 4️⃣ For non-JSON responses, return as-is
    return response

# Translation Test Endpoint 
@app.get("/test-translation")
async def test_translation(language: str = Query("en", description="Target language for translation")):
    return {
        "message": "Hello, welcome to our platform!",
        "info": "This is a sample response",
        "items": ["Apple", "Banana", "Cherry"],
        "language_received": language  # for debugging
    }

# -------------------- HEALTH CHECK --------------------
@app.get("/health")
def health_check():
    return {
        "status": "ok",
        "service": "MargDarshak Backend",
        "timestamp": utc_now()
    }
# -------------------- RUN APP --------------------
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)